from __future__ import annotations

import base64
import csv
import html
import io
import json
import mimetypes
import os
import re
import time
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import requests

try:
    import psutil
except Exception:
    psutil = None


APP_NAME = "llm Assistant"


HOME_ICON = "🤖"
ORGANIZE_ICON = "🗂️"
TRANSLATE_ICON = "🌍"
TEXT_TOOLS_ICON = "📝"
SPEED_ICON = "⚡"
SETTINGS_ICON = "🧠"
BACK_ICON = "←"

DEFAULT_ENDPOINT = "http://127.0.0.1:1234/v1/chat/completions"
DEFAULT_TIMEOUT = 120
DEFAULT_MAX_FILES = 1000
DEFAULT_MAX_CHARS = 8000
DEFAULT_MAX_FILE_MB = 25
DEFAULT_MAX_IMAGE_MB = 8
DEFAULT_TEMP = 0.10

MOUSE_SCROLL_UNITS = 1
SPEED_PROCESS_LIMIT = 10

TEXT_EXTS = {
    ".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go",
    ".rs", ".php", ".rb", ".swift", ".kt", ".kts", ".sql", ".sh", ".bat", ".ps1", ".ini", ".cfg",
    ".conf", ".log", ".toml", ".env", ".tex", ".bib", ".dart", ".lua", ".r", ".m", ".vue",
    ".svelte", ".gradle", ".properties",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
SPREADSHEET_EXTS = {".xlsx", ".xlsm", ".csv", ".tsv"}
PRESENTATION_EXTS = {".pptx"}
ARCHIVE_EXTS = {".zip"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac"}
VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".webm", ".wmv"}

INVALID_FS_CHARS = r'<>:"/\\|?*'
WINDOWS_RESERVED = {
    "CON", "PRN", "AUX", "NUL",
    "COM1", "COM2", "COM3", "COM4", "COM5", "COM6", "COM7", "COM8", "COM9",
    "LPT1", "LPT2", "LPT3", "LPT4", "LPT5", "LPT6", "LPT7", "LPT8", "LPT9",
}

SUPPORTED_VISION_MIMES = {"image/jpeg", "image/png", "image/webp"}

FORMAT_CATEGORY_MAP = {
    ".py": "Python",
    ".js": "JavaScript",
    ".ts": "TypeScript",
    ".tsx": "TypeScript",
    ".jsx": "JavaScript",
    ".java": "Java",
    ".c": "C-Cpp",
    ".cpp": "C-Cpp",
    ".h": "C-Cpp",
    ".hpp": "C-Cpp",
    ".cs": "CSharp",
    ".php": "PHP",
    ".go": "Go",
    ".rs": "Rust",
    ".sql": "SQL",
    ".html": "HTML",
    ".htm": "HTML",
    ".css": "CSS",
    ".json": "JSON",
    ".xml": "XML",
    ".md": "Markdown",
    ".txt": "Text",
    ".pdf": "PDF",
    ".docx": "Word",
    ".pptx": "Presentation",
    ".xlsx": "Spreadsheet",
    ".xlsm": "Spreadsheet",
    ".csv": "Spreadsheet",
    ".tsv": "Spreadsheet",
    ".tex": "LaTeX",
    ".bib": "LaTeX",
    ".yaml": "YAML",
    ".yml": "YAML",
    ".toml": "Config",
    ".ini": "Config",
    ".cfg": "Config",
    ".conf": "Config",
    ".env": "Config",
    ".log": "Log",
}

ROOT_ORGANIZED_FOLDERS = {
    "Python", "JavaScript", "TypeScript", "Java", "C-Cpp", "CSharp", "PHP", "Go", "Rust",
    "SQL", "HTML", "CSS", "JSON", "XML", "Markdown", "Text", "PDF", "Word", "Presentation",
    "Spreadsheet", "LaTeX", "YAML", "Config", "Log", "Images", "Nature", "Animals",
    "Screenshots", "City", "Food", "Vehicles", "People", "Documents", "Audio", "Video",
    "Archive", "Unsorted",
}

SYSTEM_USERNAMES = {"system", "local service", "network service"}

PROTECTED_PROCESS_NAMES = {
    "system", "idle", "registry", "smss.exe", "csrss.exe", "wininit.exe", "services.exe",
    "lsass.exe", "winlogon.exe", "fontdrvhost.exe", "dwm.exe", "explorer.exe",
    "svchost.exe", "taskhostw.exe", "spoolsv.exe", "searchindexer.exe",
    "searchhost.exe", "shellexperiencehost.exe", "startmenuexperiencehost.exe",
    "runtimebroker.exe", "ctfmon.exe", "securityhealthservice.exe", "msmpeng.exe",
    "wudfhost.exe", "sihost.exe", "audiodg.exe", "conhost.exe",
    "code.exe", "pycharm64.exe", "devenv.exe", "studio64.exe",
    "notepad.exe", "notepad++.exe", "winword.exe", "excel.exe", "powerpnt.exe",
    "chrome.exe", "msedge.exe", "firefox.exe", "opera.exe", "brave.exe",
    "discord.exe", "telegram.exe", "whatsapp.exe", "teams.exe", "zoom.exe",
    "obs64.exe", "python.exe", "pythonw.exe",
}

SYSTEM_PATH_HINTS = (
    r"\windows\system32",
    r"\windows\syswow64",
    r"\windows\servicing",
    r"\windows\winsxs",
)


@dataclass
class FileInsight:
    path: Path
    kind: str
    snippet: str
    meta: Dict[str, Any] = field(default_factory=dict)
    extraction_method: str = ""
    error: str = ""


@dataclass
class Decision:
    source: Path
    category: str
    title: str
    summary: str
    confidence: float
    model_used: str
    kind: str
    extraction_method: str
    raw_output: str
    snippet: str
    target: Path
    action: str = "move"
    error: str = ""
    process_pid: Optional[int] = None
    process_name: str = ""
    process_username: str = ""


def human_size(size: int) -> str:
    units = ["B", "KB", "MB", "GB", "TB"]
    value = float(size)
    for unit in units:
        if value < 1024 or unit == units[-1]:
            return f"{value:.1f} {unit}" if unit != "B" else f"{int(value)} B"
        value /= 1024
    return f"{size} B"


def normalize_whitespace(text: str) -> str:
    text = text or ""
    text = text.replace("\x00", " ")
    return re.sub(r"\s+", " ", text).strip()


def sanitize_component(text: str, fallback: str = "Untitled", max_len: int = 90) -> str:
    text = (text or "").strip()
    if not text:
        text = fallback
    text = html.unescape(text)
    text = re.sub(r"[\r\n\t]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    text = "".join("_" if ch in INVALID_FS_CHARS else ch for ch in text)
    text = text.rstrip(" .")
    text = re.sub(r"\.+", ".", text)
    text = text[:max_len].strip() or fallback
    if text.upper() in WINDOWS_RESERVED:
        text = f"_{text}"
    return text


def safe_float(value: Any, default: float = 0.5) -> float:
    try:
        x = float(value)
    except Exception:
        return default
    return max(0.0, min(1.0, x))


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    parent = path.parent
    i = 2
    while True:
        candidate = parent / f"{stem} ({i}){suffix}"
        if not candidate.exists():
            return candidate
        i += 1


def read_text_guessing(path: Path, max_chars: int) -> str:
    raw = path.read_bytes()
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return raw.decode(enc, errors="ignore")[:max_chars]
        except Exception:
            continue
    return raw.decode("utf-8", errors="ignore")[:max_chars]


def guess_kind(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in PDF_EXTS:
        return "pdf"
    if ext in DOCX_EXTS:
        return "docx"
    if ext in SPREADSHEET_EXTS:
        return "spreadsheet"
    if ext in PRESENTATION_EXTS:
        return "presentation"
    if ext in ARCHIVE_EXTS:
        return "archive"
    if ext in AUDIO_EXTS:
        return "audio"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in TEXT_EXTS:
        return "text"

    mt, _ = mimetypes.guess_type(path.name)
    if mt:
        if mt.startswith("image/"):
            return "image"
        if mt.startswith("text/"):
            return "text"
        if mt.startswith("audio/"):
            return "audio"
        if mt.startswith("video/"):
            return "video"
    return "unknown"


def likely_vision_model_name(name: str) -> bool:
    n = name.lower()
    keys = [
        "vision", "vl", "llava", "minicpm-v", "qwen-vl", "qwen2-vl",
        "qwen2.5-vl", "gemma-3", "phi-3.5-vision", "moondream",
        "pixtral", "internvl",
    ]
    return any(k in n for k in keys)


def likely_text_model_name(name: str) -> bool:
    return not likely_vision_model_name(name)


def is_system_username(username: str) -> bool:
    return (username or "").strip().lower() in SYSTEM_USERNAMES


def is_protected_process_name(name: str) -> bool:
    return (name or "").strip().lower() in PROTECTED_PROCESS_NAMES


def is_system_process_path(exe_path: str) -> bool:
    p = (exe_path or "").strip().lower()
    return any(hint in p for hint in SYSTEM_PATH_HINTS)


def get_current_process_info() -> Dict[str, Any]:
    info = {
        "pid": os.getpid(),
        "name": "",
        "exe": "",
        "username": "",
    }

    if psutil is None:
        return info

    try:
        proc = psutil.Process(os.getpid())
        info["name"] = proc.name() or ""
    except Exception:
        pass

    try:
        proc = psutil.Process(os.getpid())
        info["exe"] = proc.exe() or ""
    except Exception:
        pass

    try:
        proc = psutil.Process(os.getpid())
        info["username"] = proc.username() or ""
    except Exception:
        pass

    return info


def is_current_app_process(pid: int, name: str = "", exe_path: str = "") -> bool:
    current = get_current_process_info()

    try:
        if int(pid) == int(current.get("pid", -1)):
            return True
    except Exception:
        pass

    current_name = (current.get("name") or "").strip().lower()
    current_exe = (current.get("exe") or "").strip().lower()

    check_name = (name or "").strip().lower()
    check_exe = (exe_path or "").strip().lower()

    if current_name and check_name and current_name == check_name:
        return True

    if current_exe and check_exe and current_exe == check_exe:
        return True

    return False


def make_process_source(name: str, pid: int) -> Path:
    return Path(f"[PROCESS] {sanitize_component(name or 'process', fallback='process', max_len=60)} [{pid}]")


def filename_needs_translation(stem: str) -> bool:
    s = normalize_whitespace(stem).strip().lower()
    if not s:
        return False

    cleaned = re.sub(r"[\W_]+", "", s)
    if not cleaned:
        return False

    if cleaned.isdigit():
        return False

    if re.fullmatch(r"(img|image|dsc|pic|photo|screenshot|scan|file)?\d+", cleaned):
        return False

    return True


def normalize_category_name(raw: str, insight: FileInsight) -> str:
    text = sanitize_component(raw or "", fallback="", max_len=50).strip()
    lowered = text.lower()

    mapping = {
        "python code": "Python",
        "py": "Python",
        "script": FORMAT_CATEGORY_MAP.get(insight.path.suffix.lower(), "Text"),
        "javascript code": "JavaScript",
        "typescript code": "TypeScript",
        "text file": "Text",
        "document": "Documents" if insight.kind == "image" else FORMAT_CATEGORY_MAP.get(insight.path.suffix.lower(), "Text"),
        "doc": "Documents" if insight.kind == "image" else FORMAT_CATEGORY_MAP.get(insight.path.suffix.lower(), "Text"),
        "image": "Images",
        "photo": "Images",
        "picture": "Images",
        "screenshot": "Screenshots",
        "screen": "Screenshots",
        "ui": "Screenshots",
        "animal": "Animals",
        "vehicle": "Vehicles",
        "people": "People",
        "person": "People",
        "cityscape": "City",
        "urban": "City",
        "nature scene": "Nature",
        "natural": "Nature",
        "archive file": "Archive",
        "zip": "Archive",
        "audio file": "Audio",
        "video file": "Video",
        "excel": "Spreadsheet",
        "sheet": "Spreadsheet",
        "word": "Word",
        "powerpoint": "Presentation",
    }

    if lowered in mapping:
        return mapping[lowered]

    if not text or lowered in {"unknown", "other", "misc", "general", "file", "files", "untitled", "none", "null"}:
        ext = insight.path.suffix.lower()
        if insight.kind == "image":
            return "Images"
        if insight.kind == "audio":
            return "Audio"
        if insight.kind == "video":
            return "Video"
        if insight.kind == "archive":
            return "Archive"
        return FORMAT_CATEGORY_MAP.get(ext, "Unsorted")

    return text


class Extractor:
    @staticmethod
    def file_meta(path: Path) -> Dict[str, Any]:
        stat = path.stat()
        return {
            "name": path.name,
            "stem": path.stem,
            "suffix": path.suffix.lower(),
            "size_bytes": stat.st_size,
            "size_human": human_size(stat.st_size),
            "modified": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(stat.st_mtime)),
        }

    @classmethod
    def extract(cls, path: Path, max_chars: int = DEFAULT_MAX_CHARS) -> FileInsight:
        meta = cls.file_meta(path)
        kind = guess_kind(path)

        try:
            if kind == "text":
                return FileInsight(path, kind, cls.extract_text_like(path, max_chars), meta, "plain_text")
            if kind == "pdf":
                return FileInsight(path, kind, cls.extract_pdf(path, max_chars), meta, "pdf_text")
            if kind == "docx":
                return FileInsight(path, kind, cls.extract_docx(path, max_chars), meta, "docx_text")
            if kind == "spreadsheet":
                return FileInsight(path, kind, cls.extract_spreadsheet(path, max_chars), meta, "sheet_preview")
            if kind == "presentation":
                return FileInsight(path, kind, cls.extract_pptx(path, max_chars), meta, "pptx_text")
            if kind == "archive":
                return FileInsight(path, kind, cls.extract_zip(path, max_chars), meta, "zip_listing")
            if kind == "audio":
                return FileInsight(path, kind, cls.extract_media_stub(path, "audio", max_chars), meta, "audio_metadata")
            if kind == "video":
                return FileInsight(path, kind, cls.extract_media_stub(path, "video", max_chars), meta, "video_metadata")
            if kind == "image":
                return FileInsight(path, kind, cls.extract_image_stub(path, max_chars), meta, "image_metadata")
            return FileInsight(path, kind, cls.extract_unknown_stub(path, max_chars), meta, "fallback")
        except Exception as exc:
            return FileInsight(path, kind, "", meta, "error", error=str(exc))

    @staticmethod
    def extract_text_like(path: Path, max_chars: int) -> str:
        ext = path.suffix.lower()

        if ext in {".html", ".htm", ".xml"}:
            txt = read_text_guessing(path, max_chars * 3)
            txt = re.sub(r"<script.*?</script>", " ", txt, flags=re.I | re.S)
            txt = re.sub(r"<style.*?</style>", " ", txt, flags=re.I | re.S)
            txt = re.sub(r"<[^>]+>", " ", txt)
            txt = html.unescape(txt)
            return normalize_whitespace(txt)[:max_chars]

        if ext == ".json":
            text = read_text_guessing(path, max_chars * 2)
            try:
                obj = json.loads(text)
                return json.dumps(obj, ensure_ascii=False, indent=2)[:max_chars]
            except Exception:
                return text[:max_chars]

        return read_text_guessing(path, max_chars)

    @staticmethod
    def extract_pdf(path: Path, max_chars: int) -> str:
        text = ""
        errors: List[str] = []

        try:
            from pypdf import PdfReader  # type: ignore
            reader = PdfReader(str(path))
            for page in reader.pages[:20]:
                text += (page.extract_text() or "") + "\n"
                if len(text) >= max_chars:
                    break
        except Exception as exc:
            errors.append(f"pypdf: {exc}")

        if not text:
            try:
                import PyPDF2  # type: ignore
                with path.open("rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages[:20]:
                        text += (page.extract_text() or "") + "\n"
                        if len(text) >= max_chars:
                            break
            except Exception as exc:
                errors.append(f"PyPDF2: {exc}")

        if text.strip():
            return text[:max_chars]

        raise RuntimeError("PDF text extraction failed: " + " | ".join(errors))

    @staticmethod
    def extract_docx(path: Path, max_chars: int) -> str:
        from docx import Document  # type: ignore
        doc = Document(str(path))
        chunks: List[str] = []

        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                chunks.append(t)
            if sum(len(x) for x in chunks) >= max_chars:
                break

        if not chunks:
            for table in doc.tables:
                for row in table.rows:
                    row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_txt:
                        chunks.append(row_txt)
                    if sum(len(x) for x in chunks) >= max_chars:
                        break

        return "\n".join(chunks)[:max_chars]

    @staticmethod
    def extract_spreadsheet(path: Path, max_chars: int) -> str:
        ext = path.suffix.lower()
        if ext in {".csv", ".tsv"}:
            delimiter = "," if ext == ".csv" else "\t"
            rows: List[str] = []
            with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
                reader = csv.reader(f, delimiter=delimiter)
                for i, row in enumerate(reader):
                    rows.append(" | ".join(cell.strip() for cell in row[:12]))
                    if i >= 30:
                        break
            return "\n".join(rows)[:max_chars]

        from openpyxl import load_workbook  # type: ignore
        wb = load_workbook(filename=str(path), read_only=True, data_only=True)
        chunks: List[str] = []

        for ws in wb.worksheets[:5]:
            chunks.append(f"[Sheet] {ws.title}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                values = ["" if v is None else str(v) for v in row[:12]]
                line = " | ".join(v.strip() for v in values if v is not None and str(v).strip())
                if line:
                    chunks.append(line)
                if i >= 25 or sum(len(x) for x in chunks) >= max_chars:
                    break
            if sum(len(x) for x in chunks) >= max_chars:
                break
        return "\n".join(chunks)[:max_chars]

    @staticmethod
    def extract_pptx(path: Path, max_chars: int) -> str:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        chunks: List[str] = []

        for i, slide in enumerate(prs.slides[:30], start=1):
            chunks.append(f"[Slide {i}]")
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    chunks.append(str(text).strip())
                if sum(len(x) for x in chunks) >= max_chars:
                    break
            if sum(len(x) for x in chunks) >= max_chars:
                break
        return "\n".join(chunks)[:max_chars]

    @staticmethod
    def extract_zip(path: Path, max_chars: int) -> str:
        chunks: List[str] = []
        with zipfile.ZipFile(path, "r") as zf:
            for info in zf.infolist()[:120]:
                chunks.append(f"{info.filename} ({human_size(info.file_size)})")
                if sum(len(x) for x in chunks) >= max_chars:
                    break
        return "\n".join(chunks)[:max_chars]

    @staticmethod
    def extract_media_stub(path: Path, kind: str, max_chars: int) -> str:
        return (
            f"Filename: {path.name}\n"
            f"Type: {kind}\n"
            "This is a media file. Use filename, extension, and metadata to infer a useful broad title and category."
        )[:max_chars]

    @staticmethod
    def extract_image_stub(path: Path, max_chars: int) -> str:
        return (
            f"Filename: {path.name}\n"
            "Type: image\n"
            "Use the vision model directly for image classification."
        )[:max_chars]

    @staticmethod
    def extract_unknown_stub(path: Path, max_chars: int) -> str:
        return (
            f"Filename: {path.name}\n"
            "Type: unknown binary/unsupported\n"
            "Infer a broad category from filename, extension, and metadata only."
        )[:max_chars]


class LLMClient:
    def __init__(self, endpoint: str, timeout: int = DEFAULT_TIMEOUT):
        self.endpoint = endpoint.strip()
        self.timeout = timeout

    def _base_url(self) -> str:
        base = self.endpoint.rstrip("/")
        if base.endswith("/chat/completions"):
            base = base[:-len("/chat/completions")]
        return base

    def _post(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        response = requests.post(self.endpoint, json=payload, timeout=self.timeout)
        response.raise_for_status()
        return response.json()

    def list_models(self) -> List[str]:
        r = requests.get(self._base_url() + "/models", timeout=20)
        r.raise_for_status()
        data = r.json()
        return sorted([item.get("id", "") for item in data.get("data", []) if item.get("id")])

    @staticmethod
    def strip_reasoning(text: str) -> str:
        text = text or ""
        text = re.sub(r"<think>.*?</think>", "", text, flags=re.S | re.I)
        return text.strip()

    @staticmethod
    def extract_content(response: Dict[str, Any]) -> str:
        choices = response.get("choices") or []
        if not choices:
            return ""

        message = choices[0].get("message") or {}
        content = message.get("content", "")

        if isinstance(content, list):
            parts = []
            for item in content:
                if isinstance(item, dict) and item.get("type") == "text":
                    parts.append(item.get("text", ""))
                elif isinstance(item, str):
                    parts.append(item)
            content = "\n".join(parts)

        return LLMClient.strip_reasoning(str(content))

    @staticmethod
    def parse_jsonish(text: str) -> Dict[str, Any]:
        cleaned = (text or "").strip()
        try:
            data = json.loads(cleaned)
            if isinstance(data, dict):
                return data
        except Exception:
            pass

        match = re.search(r"\{.*\}", cleaned, flags=re.S)
        if match:
            try:
                data = json.loads(match.group(0))
                if isinstance(data, dict):
                    return data
            except Exception:
                pass

        return {
            "title": "",
            "category": "",
            "summary": "",
            "confidence": 0.5,
        }

    def _image_to_data_url(self, path: Path) -> str:
        mime = (mimetypes.guess_type(path.name)[0] or "").lower()

        if mime in SUPPORTED_VISION_MIMES:
            raw = path.read_bytes()
        else:
            try:
                from PIL import Image  # type: ignore
            except Exception as exc:
                raise RuntimeError(
                    f"Unsupported image format for vision: {path.suffix}. "
                    "Use PNG/JPG/WebP, or install Pillow to auto-convert BMP/GIF/TIFF."
                ) from exc

            with Image.open(path) as im:
                buf = io.BytesIO()
                if im.mode not in ("RGB", "RGBA"):
                    im = im.convert("RGBA" if "A" in im.mode else "RGB")
                im.save(buf, format="PNG")
                raw = buf.getvalue()
                mime = "image/png"

        b64 = base64.b64encode(raw).decode("ascii")
        return f"data:{mime};base64,{b64}"

    def _classification_payload(self, insight: FileInsight, model: str, structured: bool = True) -> Dict[str, Any]:
        system = (
            "You are a smart file organizer.\n"
            "Return ONLY valid JSON with keys: title, category, summary, confidence.\n\n"
            "Rules:\n"
            "- title = the real subject or purpose of the file.\n"
            "- category = the file format, language, or broad file family.\n"
            "- summary = one short sentence.\n"
            "- confidence = number from 0 to 1.\n\n"
            "Category rules:\n"
            "- .py -> Python\n"
            "- .js -> JavaScript\n"
            "- .ts/.tsx -> TypeScript\n"
            "- .java -> Java\n"
            "- .cpp/.c/.h/.hpp -> C-Cpp\n"
            "- .cs -> CSharp\n"
            "- .php -> PHP\n"
            "- .go -> Go\n"
            "- .rs -> Rust\n"
            "- .sql -> SQL\n"
            "- .html/.htm -> HTML\n"
            "- .css -> CSS\n"
            "- .json -> JSON\n"
            "- .xml -> XML\n"
            "- .md -> Markdown\n"
            "- .txt -> Text\n"
            "- .pdf -> PDF\n"
            "- .docx -> Word\n"
            "- .pptx -> Presentation\n"
            "- .xlsx/.csv/.tsv -> Spreadsheet\n"
            "- .tex/.bib -> LaTeX\n"
            "- screenshots -> Screenshots\n"
            "- use the extension and content together.\n\n"
            "Do not wrap JSON in markdown."
        )

        user = {
            "file_name": insight.path.name,
            "file_type": insight.kind,
            "extension": insight.path.suffix.lower(),
            "extraction_method": insight.extraction_method,
            "metadata": insight.meta,
            "content_excerpt": insight.snippet,
            "extract_error": insight.error,
        }

        payload: Dict[str, Any] = {
            "model": model,
            "temperature": DEFAULT_TEMP,
            "max_tokens": 220,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": json.dumps(user, ensure_ascii=False)},
            ],
        }

        if structured:
            payload["response_format"] = {
                "type": "json_schema",
                "json_schema": {
                    "name": "organizer_result",
                    "strict": True,
                    "schema": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "category": {"type": "string"},
                            "summary": {"type": "string"},
                            "confidence": {"type": "number"},
                        },
                        "required": ["title", "category", "summary", "confidence"],
                        "additionalProperties": False,
                    },
                },
            }

        return payload

    def classify_text(self, insight: FileInsight, model: str) -> Tuple[Dict[str, Any], str]:
        first_error = None
        try:
            data = self._post(self._classification_payload(insight, model, structured=True))
            raw = self.extract_content(data)
            return self.parse_jsonish(raw), raw
        except Exception as exc:
            first_error = exc

        data = self._post(self._classification_payload(insight, model, structured=False))
        raw = self.extract_content(data)
        parsed = self.parse_jsonish(raw)
        if not raw and first_error:
            raise first_error
        return parsed, raw

    def _image_classification_payload(self, insight: FileInsight, model: str, structured: bool = True) -> Dict[str, Any]:
        data_url = self._image_to_data_url(insight.path)
        system = (
            "You are a smart file organizer for images.\n"
            "Look at the image and return ONLY valid JSON with keys: title, image/category, summary, confidence.\n\n"
            "Rules:\n"
            "- title = a short specific image title.\n"
            "- category = a setase image/ with broad visual category.\n"
            "- summary = one short sentence describing the image.\n"
            "- confidence = number from 0 to 1.\n\n"
            "Category examples:\n"
            "- beach/mountain/forest/river/sea -> image/Nature\n"
            "- dog/cat/bird/horse/wildlife -> image/Animals\n"
            "- screenshot/UI/app/window -> image/Screenshots\n"
            "- city/buildings/street/urban -> image/City\n"
            "- food/meals/drinks -> image/Food\n"
            "- car/bike/motorcycle/truck -> image/Vehicles\n"
            "- people/portrait/selfie -> image/People\n"
            "- document/scan/paper/page -> image/Documents\n\n"
            "Do not wrap JSON in markdown."
        )

        payload: Dict[str, Any] = {
            "model": model,
            "temperature": 0.0,
            "max_tokens": 220,
            "messages": [
                {"role": "system", "content": system},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": f"Analyze this image and return JSON only.\nFilename: {insight.path.name}"},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                },
            ],
        }

        if structured:
            payload["response_format"] = {
                "type": "json_schema",
                "json_schema": {
                    "name": "image_organizer_result",
                    "strict": True,
                    "schema": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "category": {"type": "string"},
                            "summary": {"type": "string"},
                            "confidence": {"type": "number"},
                        },
                        "required": ["title", "category", "summary", "confidence"],
                        "additionalProperties": False,
                    },
                },
            }

        return payload

    def classify_image(self, insight: FileInsight, model: str) -> Tuple[Dict[str, Any], str]:
        first_error = None
        try:
            data = self._post(self._image_classification_payload(insight, model, structured=True))
            raw = self.extract_content(data)
            return self.parse_jsonish(raw), raw
        except Exception as exc:
            first_error = exc

        data = self._post(self._image_classification_payload(insight, model, structured=False))
        raw = self.extract_content(data)
        parsed = self.parse_jsonish(raw)
        if not raw and first_error:
            raise first_error
        return parsed, raw

    def translate_filename(self, insight: FileInsight, model: str, target_language: str) -> Tuple[str, str]:
        system = (
            "You translate file names.\n"
            "Return ONLY valid JSON with keys: title, summary, confidence.\n"
            "- title must be the translated meaning of the filename.\n"
            "- use the current filename text as the source.\n"
            "- do not output generic placeholders like 'Translated File Name'.\n"
            "- do not output labels, explanations, or markdown.\n"
            "- do not include file extension.\n"
            "- if the filename is only numbers or meaningless, keep it close to the original.\n"
            "- title must be short and natural.\n"
        )

        user = {
            "target_language": target_language,
            "current_filename": insight.path.stem,
            "extension": insight.path.suffix.lower(),
            "kind": insight.kind,
            "task": f"Translate this filename into {target_language}. Return JSON only.",
        }

        payload = {
            "model": model,
            "temperature": 0.1,
            "max_tokens": 120,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": json.dumps(user, ensure_ascii=False)},
            ],
            "response_format": {
                "type": "json_schema",
                "json_schema": {
                    "name": "translate_filename_result",
                    "strict": True,
                    "schema": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "summary": {"type": "string"},
                            "confidence": {"type": "number"},
                        },
                        "required": ["title", "summary", "confidence"],
                        "additionalProperties": False,
                    },
                },
            },
        }

        try:
            data = self._post(payload)
            raw = self.extract_content(data)
            parsed = self.parse_jsonish(raw)
        except Exception:
            payload.pop("response_format", None)
            data = self._post(payload)
            raw = self.extract_content(data)
            parsed = self.parse_jsonish(raw)

        title = sanitize_component(parsed.get("title") or insight.path.stem, fallback=insight.path.stem, max_len=90)

        bad_titles = {
            "translated file name",
            "translated filename",
            "file name",
            "filename",
            "translated name",
            "name translated",
            "nom du fichier traduit",
            "nom traduit",
            "اسم الملف المترجم",
            "اسم مترجم",
        }

        if title.strip().lower() in bad_titles:
            title = sanitize_component(insight.path.stem, fallback=insight.path.stem, max_len=90)

        return title, raw

    def process_text_tool(self, text: str, model: str, operation: str, target_language: str = "") -> str:
        op = (operation or "").strip().lower()
        target_language = (target_language or "").strip()

        if op == "translate":
            system = (
                "You are a professional translator.\n"
                f"Translate the user's text into {target_language or 'English'}.\n"
                "Return only the translated text.\n"
                "Do not add explanations, labels, quotes, or markdown.\n"
                "Preserve paragraphs, lists, and line breaks when possible.\n"
            )
            user_task = f"Translate this text into {target_language or 'English'}."
            temperature = 0.1
            max_tokens = 1800
        elif op == "correct grammar":
            system = (
                "You correct grammar, spelling, punctuation, and wording mistakes.\n"
                "Return only the corrected text.\n"
                "Do not explain the changes.\n"
                "Preserve the original meaning, tone, and formatting whenever possible.\n"
            )
            user_task = "Correct the grammar of this text."
            temperature = 0.1
            max_tokens = 1800
        elif op == "summarize":
            system = (
                "You create concise, accurate summaries.\n"
                "Return only the summary text.\n"
                "Do not add labels, commentary, or markdown.\n"
            )
            user_task = "Summarize this text clearly and briefly."
            temperature = 0.2
            max_tokens = 900
        else:
            raise ValueError(f"Unsupported text tool operation: {operation}")

        payload = {
            "model": model,
            "temperature": temperature,
            "max_tokens": max_tokens,
            "messages": [
                {"role": "system", "content": system},
                {
                    "role": "user",
                    "content": json.dumps({"task": user_task, "text": text}, ensure_ascii=False),
                },
            ],
        }

        data = self._post(payload)
        return self.extract_content(data)

    def classify_processes(self, processes: List[Dict[str, Any]], model: str) -> Tuple[List[Dict[str, Any]], str]:
        if not processes:
            return [], "[]"

        system = (
            "You are helping optimize a Windows computer.\n"
            "Choose which USER processes can be safely closed to recover CPU and RAM.\n"
            "Be conservative.\n"
            "Use ONLY exact process names that already exist in the provided list. Copy them exactly.\n"
            "Never invent names. Never return system or protected processes.\n"
            "Avoid browsers, IDEs, editors, Office apps, chat apps, shell, audio, security, drivers, Windows services, and anything that may contain unsaved user work unless the data strongly proves it is safe.\n"
            "Prefer background helpers, launchers, duplicate updaters, crash handlers, and non-critical heavy user processes.\n"
            "Return ONLY valid JSON in this format:\n"
            '{"items":[{"process_name":"steamwebhelper.exe","close":true,"reason":"background helper using a lot of RAM","confidence":0.87}]}'
            "\nRules:\n"
            "- process_name must exactly match one name from the input list\n"
            "- close must be true or false\n"
            "- reason must be short\n"
            "- confidence must be 0..1\n"
            "- return only the processes that should be considered, ordered from best to worst candidate\n"
        )

        user = {
            "task": "Choose which user-space processes should be closed first for safe CPU/RAM recovery.",
            "processes": processes,
        }

        payload = {
            "model": model,
            "temperature": 0.0,
            "max_tokens": 1200,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": json.dumps(user, ensure_ascii=False)},
            ],
            "response_format": {
                "type": "json_schema",
                "json_schema": {
                    "name": "process_close_plan",
                    "strict": True,
                    "schema": {
                        "type": "object",
                        "properties": {
                            "items": {
                                "type": "array",
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "process_name": {"type": "string"},
                                        "close": {"type": "boolean"},
                                        "reason": {"type": "string"},
                                        "confidence": {"type": "number"},
                                    },
                                    "required": ["process_name", "close", "reason", "confidence"],
                                    "additionalProperties": False,
                                },
                            },
                        },
                        "required": ["items"],
                        "additionalProperties": False,
                    },
                },
            },
        }

        try:
            data = self._post(payload)
            raw = self.extract_content(data).strip()
        except Exception:
            payload.pop("response_format", None)
            data = self._post(payload)
            raw = self.extract_content(data).strip()

        if not raw:
            return [], raw

        try:
            obj = json.loads(raw)
        except Exception:
            match = re.search(r"\{.*\}", raw, flags=re.S)
            if match:
                try:
                    obj = json.loads(match.group(0))
                except Exception:
                    return [], raw
            else:
                match = re.search(r"\[.*\]", raw, flags=re.S)
                if match:
                    try:
                        obj = json.loads(match.group(0))
                    except Exception:
                        return [], raw
                else:
                    return [], raw

        if isinstance(obj, dict) and isinstance(obj.get("items"), list):
            return obj["items"], raw

        if isinstance(obj, list):
            return obj, raw

        return [], raw


def build_decision(
    insight: FileInsight,
    parsed: Dict[str, Any],
    raw_output: str,
    root: Path,
    rename_files: bool,
    action: str,
    model_used: str,
) -> Decision:
    ext = insight.path.suffix.lower()
    fallback_category = FORMAT_CATEGORY_MAP.get(ext, "Unsorted")

    if insight.kind == "image":
        fallback_category = "Images"
    elif insight.kind == "audio":
        fallback_category = "Audio"
    elif insight.kind == "video":
        fallback_category = "Video"
    elif insight.kind == "archive":
        fallback_category = "Archive"

    category = normalize_category_name(str(parsed.get("category") or fallback_category), insight)
    category = sanitize_component(category, fallback=fallback_category, max_len=50)

    title_guess = parsed.get("title") or insight.path.stem
    title = sanitize_component(str(title_guess), fallback=insight.path.stem, max_len=90)
    if title.lower() in {"untitled", "unknown", "file", "image", "document"}:
        title = sanitize_component(insight.path.stem, fallback=insight.path.stem, max_len=90)

    summary = str(parsed.get("summary") or "")
    confidence = safe_float(parsed.get("confidence", 0.5), 0.5)

    filename = f"{title}{insight.path.suffix}" if rename_files else insight.path.name
    target = root / category / filename

    return Decision(
        source=insight.path,
        category=category,
        title=title,
        summary=summary,
        confidence=confidence,
        model_used=model_used,
        kind=insight.kind,
        extraction_method=insight.extraction_method,
        raw_output=raw_output,
        snippet=insight.snippet,
        target=target,
        action=action,
        error=insight.error,
    )


